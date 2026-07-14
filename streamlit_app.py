from __future__ import annotations

from datetime import datetime
from datetime import date
from html import escape
from io import BytesIO, StringIO
from pathlib import Path
import base64
import csv
import hashlib
import hmac
import json
import os
import uuid
from urllib.parse import urlsplit, urlunsplit

import streamlit as st
from sqlalchemy import create_engine, text
from sqlalchemy.exc import SQLAlchemyError
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.lib.utils import ImageReader
from reportlab.platypus import (
    Image as PdfImage,
    KeepTogether,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)


RUNTIME_DIR = Path(".streamlit_runtime")
DB_PATH = RUNTIME_DIR / "validador.db"
LOGO_PATH = Path("logo_Valenet.png")
PBKDF2_ITERATIONS = 390_000
MAX_IMAGE_SIZE_MB = 5
_DATABASE_ENGINE = None
_DATABASE_INITIALIZED = False
APP_BUILD = "admin-recovery-bootstrap-2026-07-14-1"
RECOVERY_ADMIN_USERNAME = "walace_admin"
RECOVERY_ADMIN_PASSWORD_HASH = "6150bb431cbe4a22b3fa8e27ea180801$fe467528a5ac71d4e57e4c8732f7b3ba3ea7440a069efa9445ed645673df02c6"
RECOVERY_ADMIN_VERSION = "2026-07-14-admin-recovery-1"


st.set_page_config(
    page_title="Validador WhatsApp",
    page_icon="✅",
    layout="wide",
)


def get_credentials() -> tuple[str, str]:
    username = os.getenv("APP_USERNAME") or get_secret("APP_USERNAME", "")
    password = os.getenv("APP_PASSWORD") or get_secret("APP_PASSWORD", "")
    return username, password


def get_admin_credentials() -> tuple[str, str]:
    username = os.getenv("APP_ADMIN_USERNAME") or get_secret("APP_ADMIN_USERNAME", "Admin")
    password = os.getenv("APP_ADMIN_PASSWORD") or get_secret("APP_ADMIN_PASSWORD", "ValenetAdmin2026")
    return username, password


def get_reset_request() -> tuple[str, str, str]:
    username = os.getenv("APP_RESET_USERNAME") or get_secret("APP_RESET_USERNAME", "")
    password = os.getenv("APP_RESET_PASSWORD") or get_secret("APP_RESET_PASSWORD", "")
    version = os.getenv("APP_RESET_VERSION") or get_secret("APP_RESET_VERSION", "")
    return normalize_username(username), password, version


def get_secret(key: str, default: str = "") -> str:
    try:
        return st.secrets.get(key, default)
    except Exception:
        return default


def get_initial_password_for_user(username: str) -> str:
    app_username, app_password = get_credentials()
    admin_username, admin_password = get_admin_credentials()
    reset_username, reset_password, _ = get_reset_request()
    if username == reset_username:
        return reset_password
    if username == admin_username:
        return admin_password
    if username == app_username:
        return app_password
    return ""


def get_database_url() -> str:
    database_url = os.getenv("DATABASE_URL") or get_secret("DATABASE_URL", "")
    if database_url.startswith("postgres://"):
        database_url = database_url.replace("postgres://", "postgresql+psycopg2://", 1)
    elif database_url.startswith("postgresql://"):
        database_url = database_url.replace("postgresql://", "postgresql+psycopg2://", 1)

    if database_url:
        return database_url

    RUNTIME_DIR.mkdir(parents=True, exist_ok=True)
    return f"sqlite:///{DB_PATH.resolve().as_posix()}"


def safe_database_label(database_url: str) -> str:
    if not database_url.startswith("postgres"):
        return str(DB_PATH.resolve())

    parsed = urlsplit(database_url)
    host = parsed.hostname or "postgres"
    database = parsed.path.lstrip("/") or "database"
    return urlunsplit((parsed.scheme, host, f"/{database}", "", ""))


def get_engine():
    global _DATABASE_ENGINE
    if _DATABASE_ENGINE is None:
        _DATABASE_ENGINE = create_engine(
            get_database_url(),
            future=True,
            pool_pre_ping=True,
            pool_recycle=300,
        )
    return _DATABASE_ENGINE


def execute_statement(sql: str, params: dict[str, object] | None = None) -> None:
    with get_engine().begin() as connection:
        connection.execute(text(sql), params or {})


def fetch_one(sql: str, params: dict[str, object] | None = None):
    with get_engine().connect() as connection:
        return connection.execute(text(sql), params or {}).mappings().fetchone()


def fetch_all(sql: str, params: dict[str, object] | None = None) -> list[dict[str, object]]:
    with get_engine().connect() as connection:
        return list(connection.execute(text(sql), params or {}).mappings().fetchall())


def hash_password(password: str, salt: str | None = None) -> str:
    salt = salt or uuid.uuid4().hex
    digest = hashlib.pbkdf2_hmac(
        "sha256",
        password.encode("utf-8"),
        salt.encode("utf-8"),
        PBKDF2_ITERATIONS,
    ).hex()
    return f"{salt}${digest}"


def verify_password(password: str, password_hash: str) -> bool:
    try:
        salt, expected_digest = password_hash.split("$", 1)
    except ValueError:
        return False

    candidate = hash_password(password, salt).split("$", 1)[1]
    return hmac.compare_digest(candidate, expected_digest)


def initialize_database() -> None:
    global _DATABASE_INITIALIZED
    if _DATABASE_INITIALIZED:
        return

    expected_username, initial_password = get_credentials()
    admin_username, admin_password = get_admin_credentials()
    execute_statement(
        """
        CREATE TABLE IF NOT EXISTS users (
            username TEXT PRIMARY KEY,
            password_hash TEXT NOT NULL,
            must_change_password INTEGER NOT NULL DEFAULT 1,
            role TEXT NOT NULL DEFAULT 'user',
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL,
            last_login_at TEXT
        )
        """
    )
    execute_statement(
        """
        CREATE TABLE IF NOT EXISTS saved_reports (
            id TEXT PRIMARY KEY,
            username TEXT NOT NULL,
            audit_name TEXT,
            channel TEXT,
            auditor TEXT,
            audit_date TEXT,
            tests_json TEXT NOT NULL,
            created_at TEXT NOT NULL,
            FOREIGN KEY(username) REFERENCES users(username)
        )
        """
    )
    execute_statement(
        """
        CREATE TABLE IF NOT EXISTS app_settings (
            key TEXT PRIMARY KEY,
            value TEXT NOT NULL,
            updated_at TEXT NOT NULL
        )
        """
    )

    for alter_sql in (
        "ALTER TABLE users ADD COLUMN role TEXT NOT NULL DEFAULT 'user'",
        "ALTER TABLE users ADD COLUMN last_login_at TEXT",
    ):
        try:
            execute_statement(alter_sql)
        except SQLAlchemyError:
            pass

    create_default_user(expected_username, initial_password, "user")
    create_default_user(admin_username, admin_password, "admin")
    apply_secret_password_reset()
    apply_emergency_admin_recovery()
    _DATABASE_INITIALIZED = True


def create_default_user(username: str, password: str, role: str) -> None:
    if not username or not password:
        return

    existing = fetch_one(
        "SELECT username FROM users WHERE username = :username",
        {"username": username},
    )
    if existing is not None:
        return

    existing_role = fetch_one(
        "SELECT username FROM users WHERE role = :role LIMIT 1",
        {"role": role},
    )
    if existing_role is not None:
        return

    now = datetime.now().isoformat(timespec="seconds")
    execute_statement(
        """
        INSERT INTO users (
            username, password_hash, must_change_password, role, created_at, updated_at
        )
        VALUES (:username, :password_hash, 1, :role, :created_at, :updated_at)
        """,
        {
            "username": username,
            "password_hash": hash_password(password),
            "role": role,
            "created_at": now,
            "updated_at": now,
        },
    )


def create_user(username: str, password: str, role: str = "user", must_change_password: bool = True) -> None:
    username = normalize_username(username)
    now = datetime.now().isoformat(timespec="seconds")
    execute_statement(
        """
        INSERT INTO users (
            username, password_hash, must_change_password, role, created_at, updated_at
        )
        VALUES (
            :username, :password_hash, :must_change_password, :role, :created_at, :updated_at
        )
        """,
        {
            "username": username,
            "password_hash": hash_password(password),
            "must_change_password": 1 if must_change_password else 0,
            "role": role,
            "created_at": now,
            "updated_at": now,
        },
    )


def get_user(username: str):
    initialize_database()
    return fetch_one(
        "SELECT * FROM users WHERE username = :username",
        {"username": username},
    )


def save_new_password(username: str, password: str) -> None:
    initialize_database()
    now = datetime.now().isoformat(timespec="seconds")
    execute_statement(
        """
        UPDATE users
        SET password_hash = :password_hash, must_change_password = 0, updated_at = :updated_at
        WHERE username = :username
        """,
        {
            "password_hash": hash_password(password),
            "updated_at": now,
            "username": username,
        },
    )


def reset_user_password(username: str, temporary_password: str) -> None:
    initialize_database()
    now = datetime.now().isoformat(timespec="seconds")
    execute_statement(
        """
        UPDATE users
        SET password_hash = :password_hash,
            must_change_password = 1,
            updated_at = :updated_at
        WHERE username = :username
        """,
        {
            "password_hash": hash_password(temporary_password),
            "updated_at": now,
            "username": username,
        },
    )


def apply_secret_password_reset() -> None:
    username, temporary_password, reset_version = get_reset_request()
    if not username or not temporary_password or not reset_version:
        return

    setting_key = f"password_reset:{username}"
    applied = fetch_one(
        "SELECT value FROM app_settings WHERE key = :key",
        {"key": setting_key},
    )
    if applied and applied["value"] == reset_version:
        return

    force_secret_password_reset(username, temporary_password, reset_version)


def apply_emergency_admin_recovery() -> None:
    setting_key = f"admin_recovery:{RECOVERY_ADMIN_USERNAME}"
    applied = fetch_one(
        "SELECT value FROM app_settings WHERE key = :key",
        {"key": setting_key},
    )
    if applied and applied["value"] == RECOVERY_ADMIN_VERSION:
        return

    now = datetime.now().isoformat(timespec="seconds")
    execute_statement(
        """
        INSERT INTO users (
            username, password_hash, must_change_password, role, created_at, updated_at
        )
        VALUES (
            :username, :password_hash, 1, 'admin', :created_at, :updated_at
        )
        ON CONFLICT(username) DO UPDATE SET
            password_hash = excluded.password_hash,
            must_change_password = 1,
            role = 'admin',
            updated_at = excluded.updated_at
        """,
        {
            "username": RECOVERY_ADMIN_USERNAME,
            "password_hash": RECOVERY_ADMIN_PASSWORD_HASH,
            "created_at": now,
            "updated_at": now,
        },
    )
    execute_statement(
        """
        INSERT INTO app_settings (key, value, updated_at)
        VALUES (:key, :value, :updated_at)
        ON CONFLICT(key) DO UPDATE SET
            value = excluded.value,
            updated_at = excluded.updated_at
        """,
        {
            "key": setting_key,
            "value": RECOVERY_ADMIN_VERSION,
            "updated_at": now,
        },
    )


def force_secret_password_reset(username: str, temporary_password: str, reset_version: str) -> None:
    setting_key = f"password_reset:{username}"
    if get_user_without_init(username):
        reset_user_password_without_init(username, temporary_password)
    else:
        create_user(username, temporary_password, role="user", must_change_password=True)

    now = datetime.now().isoformat(timespec="seconds")
    execute_statement(
        """
        INSERT INTO app_settings (key, value, updated_at)
        VALUES (:key, :value, :updated_at)
        ON CONFLICT(key) DO UPDATE SET
            value = excluded.value,
            updated_at = excluded.updated_at
        """,
        {
            "key": setting_key,
            "value": reset_version,
            "updated_at": now,
        },
    )


def get_user_without_init(username: str):
    return fetch_one(
        "SELECT * FROM users WHERE username = :username",
        {"username": username},
    )


def reset_user_password_without_init(username: str, temporary_password: str) -> None:
    now = datetime.now().isoformat(timespec="seconds")
    execute_statement(
        """
        UPDATE users
        SET password_hash = :password_hash,
            must_change_password = 1,
            updated_at = :updated_at
        WHERE username = :username
        """,
        {
            "password_hash": hash_password(temporary_password),
            "updated_at": now,
            "username": username,
        },
    )


def delete_user(username: str) -> None:
    initialize_database()
    with get_engine().begin() as connection:
        connection.execute(
            text("DELETE FROM saved_reports WHERE username = :username"),
            {"username": username},
        )
        connection.execute(
            text("DELETE FROM users WHERE username = :username"),
            {"username": username},
        )


def normalize_username(username: str) -> str:
    return " ".join(username.strip().split())


def rename_user(old_username: str, new_username: str) -> None:
    new_username = normalize_username(new_username)
    if old_username == new_username:
        return

    now = datetime.now().isoformat(timespec="seconds")
    with get_engine().begin() as connection:
        connection.execute(
            text(
                """
                UPDATE users
                SET username = :new_username, updated_at = :updated_at
                WHERE username = :old_username
                """
            ),
            {
                "new_username": new_username,
                "old_username": old_username,
                "updated_at": now,
            },
        )
        connection.execute(
            text(
                """
                UPDATE saved_reports
                SET username = :new_username
                WHERE username = :old_username
                """
            ),
            {
                "new_username": new_username,
                "old_username": old_username,
            },
        )


def validate_login(username: str, password: str) -> tuple[bool, bool]:
    username = normalize_username(username)
    reset_username, reset_password, reset_version = get_reset_request()
    if (
        reset_username
        and reset_password
        and reset_version
        and username == reset_username
        and hmac.compare_digest(password, reset_password)
    ):
        force_secret_password_reset(reset_username, reset_password, reset_version)
        return True, True

    user = get_user(username)
    if not user:
        return False, False

    if not verify_password(password, user["password_hash"]):
        return False, False

    execute_statement(
        "UPDATE users SET last_login_at = :last_login_at WHERE username = :username",
        {
            "last_login_at": datetime.now().isoformat(timespec="seconds"),
            "username": username,
        },
    )

    if user["must_change_password"]:
        return True, True

    return True, False


def verify_current_password(password: str) -> bool:
    username = st.session_state.get("current_user")
    if not username:
        return False

    user = get_user(username)
    return bool(user and verify_password(password, user["password_hash"]))


def init_state() -> None:
    initialize_database()
    defaults = {
        "authenticated": False,
        "pending_password_change": False,
        "current_user": "",
        "current_role": "",
        "login_username": "",
        "tests": [],
        "audit_name": "",
        "channel": "",
        "auditor": "",
        "audit_date": date.today(),
        "editing_id": None,
        "current_report_id": "",
        "current_report_owner": "",
        "upload_version": 0,
        "report_saved_message": "",
        "sub_expected_count": 0,
    }

    for key, value in defaults.items():
        st.session_state.setdefault(key, value)


def login_screen() -> None:
    expected_username, expected_password = get_credentials()
    reset_username, reset_password, reset_version = get_reset_request()

    st.markdown(
        """
        <style>
          .login-card {
            max-width: 460px;
            margin: 8vh auto 0;
            padding: 2rem;
            border: 1px solid #d8e2dc;
            border-radius: 8px;
            background: #ffffff;
          }
        </style>
        """,
        unsafe_allow_html=True,
    )

    with st.container(border=True):
        st.title("Validador de automações WhatsApp")
        st.caption("Acesso restrito para validação de fluxos de atendimento.")

        with st.form("login_form"):
            username = st.text_input("Usuário")
            password = st.text_input("Senha", type="password")
            submitted = st.form_submit_button("Entrar", use_container_width=True)

        if submitted:
            valid_login, must_change_password = validate_login(username, password)
            if valid_login and must_change_password:
                st.session_state.login_username = username
                st.session_state.pending_password_change = True
                st.session_state.authenticated = False
                st.rerun()
            elif valid_login:
                user = get_user(username)
                st.session_state.current_user = username
                st.session_state.current_role = user["role"] if user else "user"
                st.session_state.login_username = ""
                st.session_state.pending_password_change = False
                st.session_state.authenticated = True
                st.rerun()
            else:
                st.error("Usuário ou senha inválidos.")
                typed_username = normalize_username(username)
                st.caption(
                    "Diagnóstico de reset: "
                    f"usuário configurado={reset_username or '(vazio)'}, "
                    f"usuário digitado={typed_username or '(vazio)'}, "
                    f"senha configurada={len(reset_password)} caracteres, "
                    f"senha digitada={len(password)} caracteres, "
                    f"versão={reset_version or '(vazia)'}."
                )

        if not expected_username or not expected_password:
            st.warning("Configure APP_USERNAME e APP_PASSWORD nos secrets do Streamlit.")

        st.caption(f"Versão: {APP_BUILD}")


def password_change_screen() -> None:
    username = st.session_state.get("login_username")
    initial_password = get_initial_password_for_user(username or "")

    with st.container(border=True):
        st.title("Primeiro acesso")
        st.caption("A senha inicial é provisória. Escolha seu usuário definitivo e cadastre uma nova senha.")

        with st.form("change_password_form"):
            new_username = st.text_input(
                "Novo nome de usuário",
                value=username or "",
                help="Este será o usuário usado nos próximos acessos.",
            )
            new_password = st.text_input("Nova senha", type="password")
            confirm_password = st.text_input("Confirmar nova senha", type="password")
            submitted = st.form_submit_button("Salvar usuário e senha", use_container_width=True)

        if not submitted:
            return

        normalized_username = normalize_username(new_username)
        if len(normalized_username) < 3:
            st.error("O novo usuário precisa ter pelo menos 3 caracteres.")
            return

        if not username:
            st.error("Sessão de troca de senha expirada. Faça login novamente.")
            st.session_state.pending_password_change = False
            return

        existing_user = get_user(normalized_username)
        if existing_user and normalized_username != username:
            st.error("Este usuário já existe. Escolha outro nome.")
            return

        if len(new_password) < 8:
            st.error("A nova senha precisa ter pelo menos 8 caracteres.")
            return

        if new_password == initial_password:
            st.error("A nova senha deve ser diferente da senha provisória.")
            return

        if new_password != confirm_password:
            st.error("A confirmação não confere com a nova senha.")
            return

        rename_user(username, normalized_username)
        save_new_password(normalized_username, new_password)
        user = get_user(normalized_username)
        st.session_state.current_user = normalized_username
        st.session_state.current_role = user["role"] if user else "user"
        st.session_state.login_username = ""
        st.session_state.pending_password_change = False
        st.session_state.authenticated = True
        st.success("Usuário e senha alterados com sucesso.")
        st.rerun()


def status_label(status: str) -> str:
    labels = {
        "conforme": "Conforme",
        "nao-conforme": "Não conforme",
        "pendente": "Pendente",
    }
    return labels.get(status, "Pendente")


def calculate_summary(tests: list[dict[str, object]] | None = None) -> dict[str, int]:
    tests = st.session_state.tests if tests is None else tests
    conform = sum(1 for item in tests if item["status"] == "conforme")
    non_conform = sum(1 for item in tests if item["status"] == "nao-conforme")
    evaluated = conform + non_conform
    rate = round((conform / evaluated) * 100) if evaluated else 0

    return {
        "total": len(tests),
        "conform": conform,
        "non_conform": non_conform,
        "rate": rate,
    }


def make_csv() -> str:
    output = StringIO()
    writer = csv.writer(output, delimiter=";")
    summary = calculate_summary()

    writer.writerow(["Bateria", st.session_state.audit_name])
    writer.writerow(["Canal", st.session_state.channel])
    writer.writerow(["Responsável", st.session_state.auditor])
    writer.writerow(["Data", st.session_state.audit_date.strftime("%d/%m/%Y")])
    writer.writerow(["Total", summary["total"]])
    writer.writerow(["Conformes", summary["conform"]])
    writer.writerow(["Não conformes", summary["non_conform"]])
    writer.writerow(["Percentual de conformidade", f'{summary["rate"]}%'])
    writer.writerow([])
    writer.writerow(
        ["Teste", "Status", "Cenário", "Esperado", "Sub respostas esperadas", "Observações", "Anexos"]
    )

    for test in st.session_state.tests:
        attachment_names = ", ".join(
            attachment["name"] for attachment in test.get("attachments", [])
        )
        writer.writerow(
            [
                test["title"],
                status_label(test["status"]),
                test["scenario"],
                test["expected"],
                " | ".join(test.get("sub_expected", [])),
                test["notes"],
                attachment_names,
            ]
        )

    return output.getvalue()


def serialize_tests(tests: list[dict[str, object]]) -> str:
    serializable_tests = []
    for test in tests:
        serializable = dict(test)
        serializable["attachments"] = [
            {
                "name": attachment["name"],
                "type": attachment["type"],
                "data": base64.b64encode(attachment["data"]).decode("ascii"),
            }
            for attachment in test.get("attachments", [])
        ]
        serializable_tests.append(serializable)
    return json.dumps(serializable_tests, ensure_ascii=False)


def deserialize_tests(tests_json: str) -> list[dict[str, object]]:
    tests = json.loads(tests_json)
    for test in tests:
        test["attachments"] = [
            {
                "name": attachment["name"],
                "type": attachment["type"],
                "data": base64.b64decode(attachment["data"]),
            }
            for attachment in test.get("attachments", [])
        ]
    return tests


def save_current_report_snapshot() -> None:
    username = st.session_state.get("current_user")
    if not username:
        st.error("Faça login novamente para salvar o relatório.")
        return

    report_id = st.session_state.get("current_report_id")
    report_owner = st.session_state.get("current_report_owner") or username
    payload = {
        "id": report_id or str(uuid.uuid4()),
        "username": username,
        "owner": report_owner,
        "audit_name": st.session_state.audit_name,
        "channel": st.session_state.channel,
        "auditor": st.session_state.auditor,
        "audit_date": st.session_state.audit_date.isoformat(),
        "tests_json": serialize_tests(st.session_state.tests),
        "created_at": datetime.now().isoformat(timespec="seconds"),
    }

    if report_id:
        existing = fetch_one(
            "SELECT id FROM saved_reports WHERE id = :id AND username = :owner",
            {"id": report_id, "owner": report_owner},
        )
        if not existing:
            st.error("Relatório original não encontrado. Nada foi salvo para evitar duplicidade.")
            return

        execute_statement(
            """
            UPDATE saved_reports
            SET audit_name = :audit_name,
                channel = :channel,
                auditor = :auditor,
                audit_date = :audit_date,
                tests_json = :tests_json,
                created_at = :created_at
            WHERE id = :id AND username = :owner
            """,
            payload,
        )
        st.session_state.report_saved_message = (
            "Relatório atualizado com sucesso. O formulário foi limpo para iniciar uma nova bateria."
        )
    else:
        execute_statement(
            """
            INSERT INTO saved_reports (
                id, username, audit_name, channel, auditor, audit_date, tests_json, created_at
            )
            VALUES (
                :id, :username, :audit_name, :channel, :auditor, :audit_date, :tests_json, :created_at
            )
            """,
            payload,
        )
        st.session_state.report_saved_message = (
            "Relatório salvo com sucesso. O formulário foi limpo para iniciar uma nova bateria."
        )
    clear_current_report()


def clear_current_report() -> None:
    st.session_state.audit_name = ""
    st.session_state.channel = ""
    st.session_state.auditor = ""
    st.session_state.audit_date = date.today()
    st.session_state.tests = []
    st.session_state.current_report_id = ""
    st.session_state.current_report_owner = ""
    reset_form()


def list_saved_reports() -> list[dict[str, object]]:
    username = st.session_state.get("current_user")
    if not username:
        return []

    return fetch_all(
        """
        SELECT *
        FROM saved_reports
        WHERE username = :username
        ORDER BY created_at DESC
        """,
        {"username": username},
    )


def load_saved_report(report_id: str) -> None:
    username = st.session_state.get("current_user")
    report = fetch_one(
        """
        SELECT *
        FROM saved_reports
        WHERE id = :id AND username = :username
        """,
        {"id": report_id, "username": username},
    )

    if not report:
        st.error("Relatório salvo não encontrado.")
        return

    st.session_state.audit_name = report["audit_name"] or ""
    st.session_state.channel = report["channel"] or ""
    st.session_state.auditor = report["auditor"] or ""
    st.session_state.audit_date = date.fromisoformat(report["audit_date"])
    st.session_state.tests = deserialize_tests(report["tests_json"])
    st.session_state.current_report_id = report["id"]
    st.session_state.current_report_owner = report["username"]
    st.session_state.report_saved_message = "Relatório carregado para edição."


def clear_saved_reports() -> None:
    username = st.session_state.get("current_user")
    execute_statement(
        "DELETE FROM saved_reports WHERE username = :username",
        {"username": username},
    )
    st.success("Relatórios salvos foram limpos do banco local.")


def delete_saved_report(report_id: str) -> None:
    username = st.session_state.get("current_user")
    report = fetch_one(
        """
        SELECT id
        FROM saved_reports
        WHERE id = :id AND username = :username
        """,
        {"id": report_id, "username": username},
    )
    if not report:
        st.error("Relatório selecionado não foi encontrado.")
        return

    execute_statement(
        """
        DELETE FROM saved_reports
        WHERE id = :id AND username = :username
        """,
        {"id": report_id, "username": username},
    )
    if st.session_state.get("current_report_id") == report_id:
        clear_current_report()
    st.success("Relatório selecionado excluído com sucesso.")


def list_all_users() -> list[dict[str, object]]:
    return fetch_all(
        """
        SELECT
            u.username,
            u.role,
            u.must_change_password,
            u.created_at,
            u.updated_at,
            u.last_login_at,
            COUNT(r.id) AS saved_reports
        FROM users u
        LEFT JOIN saved_reports r ON r.username = u.username
        GROUP BY
            u.username, u.role, u.must_change_password, u.created_at, u.updated_at, u.last_login_at
        ORDER BY u.role DESC, u.username ASC
        """
    )


def list_all_saved_reports() -> list[dict[str, object]]:
    return fetch_all(
        """
        SELECT *
        FROM saved_reports
        ORDER BY created_at DESC
        """
    )


def load_any_saved_report(report_id: str) -> None:
    report = fetch_one(
        """
        SELECT *
        FROM saved_reports
        WHERE id = :id
        """,
        {"id": report_id},
    )

    if not report:
        st.error("Relatório salvo não encontrado.")
        return

    st.session_state.audit_name = report["audit_name"] or ""
    st.session_state.channel = report["channel"] or ""
    st.session_state.auditor = report["auditor"] or ""
    st.session_state.audit_date = date.fromisoformat(report["audit_date"])
    st.session_state.tests = deserialize_tests(report["tests_json"])
    st.session_state.current_report_id = report["id"]
    st.session_state.current_report_owner = report["username"]
    st.session_state.report_saved_message = f"Relatório de {report['username']} carregado para edição."


def delete_any_saved_report(report_id: str) -> None:
    report = fetch_one(
        """
        SELECT id
        FROM saved_reports
        WHERE id = :id
        """,
        {"id": report_id},
    )
    if not report:
        st.error("Relatório selecionado não foi encontrado.")
        return

    execute_statement(
        """
        DELETE FROM saved_reports
        WHERE id = :id
        """,
        {"id": report_id},
    )
    if st.session_state.get("current_report_id") == report_id:
        clear_current_report()
    st.success("Relatório selecionado excluído pelo administrador.")


def database_status() -> dict[str, object]:
    initialize_database()
    users = fetch_one("SELECT COUNT(*) AS total FROM users")["total"]
    reports = fetch_one("SELECT COUNT(*) AS total FROM saved_reports")["total"]
    database_url = get_database_url()
    is_postgres = database_url.startswith("postgres")

    return {
        "installed": True,
        "path": safe_database_label(database_url),
        "backend": "PostgreSQL externo" if is_postgres else "SQLite local temporário",
        "persistent": is_postgres,
        "users": users,
        "reports": reports,
    }


def safe_report_base_name(raw_name: str | None) -> str:
    raw_name = raw_name or "testes-whatsapp"
    safe_name = "".join(
        char.lower() if char.isalnum() else "-"
        for char in raw_name
    ).strip("-")
    return safe_name or "testes-whatsapp"


def report_base_name() -> str:
    return safe_report_base_name(st.session_state.audit_name)


def make_saved_report_pdf(report: dict[str, object]) -> bytes:
    return make_pdf_report(
        audit_name=report["audit_name"] or "",
        channel=report["channel"] or "",
        auditor=report["auditor"] or "",
        audit_date=date.fromisoformat(report["audit_date"]),
        tests=deserialize_tests(report["tests_json"]),
    )


def executive_reading(summary: dict[str, int]) -> str:
    rate = summary["rate"]
    if summary["total"] == 0:
        return "Nenhum teste foi registrado nesta bateria."
    if rate >= 90:
        return "A bateria apresenta alta aderencia aos criterios esperados."
    if rate >= 70:
        return "A bateria apresenta aderencia parcial e requer acompanhamento dos itens nao conformes."
    return "A bateria indica pontos criticos que exigem plano de acao antes de expansao ou aceite operacional."


def pdf_text(value: object, fallback: str = "Nao informado") -> str:
    text = str(value or fallback)
    return escape(text)


def draw_pdf_footer(canvas, doc) -> None:
    canvas.saveState()
    canvas.setFont("Helvetica", 7)
    canvas.setFillColor(colors.HexColor("#6b7b83"))
    canvas.drawString(doc.leftMargin, 0.75 * cm, f"Pagina {doc.page}")

    logo_file = LOGO_PATH if LOGO_PATH.exists() else Path("../logo_Valenet.png")
    if logo_file.exists():
        try:
            image = ImageReader(str(logo_file))
            width, height = image.getSize()
            logo_width = 3.1 * cm
            logo_height = logo_width * (height / width)
            x = doc.pagesize[0] - doc.rightMargin - logo_width
            y = 0.45 * cm
            canvas.drawImage(
                image,
                x,
                y,
                width=logo_width,
                height=logo_height,
                preserveAspectRatio=True,
                mask="auto",
            )
        except Exception:
            canvas.drawRightString(
                doc.pagesize[0] - doc.rightMargin,
                0.75 * cm,
                "VALENET",
            )
    else:
        canvas.drawRightString(
            doc.pagesize[0] - doc.rightMargin,
            0.75 * cm,
            "VALENET",
        )

    canvas.restoreState()


def make_pdf_report(
    audit_name: str | None = None,
    channel: str | None = None,
    auditor: str | None = None,
    audit_date: date | None = None,
    tests: list[dict[str, object]] | None = None,
) -> bytes:
    audit_name = st.session_state.audit_name if audit_name is None else audit_name
    channel = st.session_state.channel if channel is None else channel
    auditor = st.session_state.auditor if auditor is None else auditor
    audit_date = st.session_state.audit_date if audit_date is None else audit_date
    tests = st.session_state.tests if tests is None else tests

    buffer = BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        rightMargin=1.5 * cm,
        leftMargin=1.5 * cm,
        topMargin=1.5 * cm,
        bottomMargin=2.2 * cm,
        title="Relatorio Executivo de Validacao",
    )
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "ReportTitle",
        parent=styles["Title"],
        fontSize=20,
        leading=24,
        textColor=colors.HexColor("#143642"),
        spaceAfter=14,
    )
    heading_style = ParagraphStyle(
        "ReportHeading",
        parent=styles["Heading2"],
        fontSize=13,
        leading=16,
        textColor=colors.HexColor("#143642"),
        spaceBefore=10,
        spaceAfter=8,
    )
    body_style = ParagraphStyle(
        "ReportBody",
        parent=styles["BodyText"],
        fontSize=9,
        leading=12,
        spaceAfter=5,
    )
    small_style = ParagraphStyle(
        "ReportSmall",
        parent=styles["BodyText"],
        fontSize=8,
        leading=10,
        textColor=colors.HexColor("#465a63"),
    )

    summary = calculate_summary(tests)
    pending = summary["total"] - summary["conform"] - summary["non_conform"]
    story = [
        Paragraph("Relatorio executivo de validacao", title_style),
        Paragraph("Automacoes de atendimento WhatsApp / call center", styles["Heading3"]),
        Spacer(1, 8),
        Paragraph(
            f"<b>Bateria:</b> {pdf_text(audit_name)}<br/>"
            f"<b>Canal:</b> {pdf_text(channel)}<br/>"
            f"<b>Responsavel:</b> {pdf_text(auditor)}<br/>"
            f"<b>Data:</b> {audit_date.strftime('%d/%m/%Y')}",
            body_style,
        ),
        Spacer(1, 12),
        Paragraph("Resumo executivo", heading_style),
        Paragraph(executive_reading(summary), body_style),
        Paragraph(
            "Este relatorio consolida os testes registrados no aplicativo, "
            "incluindo resultados, cenarios, respostas esperadas, observacoes e evidencias anexadas.",
            body_style,
        ),
        Spacer(1, 10),
    ]

    metrics_table = Table(
        [
            ["Total de testes", "Conformes", "Nao conformes", "Pendentes", "Conformidade"],
            [
                summary["total"],
                summary["conform"],
                summary["non_conform"],
                pending,
                f'{summary["rate"]}%',
            ],
        ],
        colWidths=[3.2 * cm, 3.2 * cm, 3.2 * cm, 3.2 * cm, 3.2 * cm],
    )
    metrics_table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#143642")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("BACKGROUND", (0, 1), (-1, 1), colors.HexColor("#eef5f3")),
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTNAME", (0, 1), (-1, 1), "Helvetica-Bold"),
                ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#b7c9c3")),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
                ("TOPPADDING", (0, 0), (-1, -1), 8),
            ]
        )
    )
    story.extend(
        [
            metrics_table,
            Spacer(1, 12),
            Paragraph("Cobertura da bateria", heading_style),
            Paragraph(
                f"Foram registrados {summary['total']} teste(s): "
                f"{summary['conform']} conforme(s), "
                f"{summary['non_conform']} nao conforme(s) e "
                f"{pending} pendente(s).",
                body_style,
            ),
            Spacer(1, 6),
            Paragraph("Detalhamento dos testes", heading_style),
        ]
    )

    if not tests:
        story.append(Paragraph("Nenhum teste cadastrado.", body_style))
    else:
        for index, test in enumerate(tests, start=1):
            header_block = [
                Paragraph(f"Teste {index}: {pdf_text(test['title'])}", heading_style),
                Paragraph(
                    f"<b>Resultado:</b> {status_label(test['status'])}",
                    body_style,
                ),
            ]
            story.append(KeepTogether(header_block))
            story.append(
                Paragraph(
                    f"<b>Cenario:</b> {pdf_text(test['scenario'], 'Nao informado.')}",
                    body_style,
                )
            )
            story.append(
                Paragraph(
                    f"<b>Resposta esperada:</b> {pdf_text(test['expected'], 'Nao informado.')}",
                    body_style,
                )
            )
            sub_expected = test.get("sub_expected", [])
            if sub_expected:
                story.append(Paragraph("<b>Sub respostas esperadas:</b>", body_style))
                for sub_index, sub_value in enumerate(sub_expected, start=1):
                    story.append(
                        Paragraph(
                            f"{sub_index}. {pdf_text(sub_value)}",
                            body_style,
                        )
                    )
            story.append(
                Paragraph(
                    f"<b>Observacoes:</b> {pdf_text(test['notes'], 'Nao informado.')}",
                    body_style,
                )
            )

            attachments = test.get("attachments", [])
            if attachments:
                story.append(Spacer(1, 6))
                story.append(Paragraph(f"Evidencias anexadas ({len(attachments)})", small_style))
                for attachment in attachments:
                    try:
                        image_reader = ImageReader(BytesIO(attachment["data"]))
                        width, height = image_reader.getSize()
                        max_width = 14 * cm
                        max_height = 7.5 * cm
                        scale = min(max_width / width, max_height / height)
                        image = PdfImage(
                            BytesIO(attachment["data"]),
                            width=width * scale,
                            height=height * scale,
                        )
                        story.append(image)
                        story.append(Paragraph(pdf_text(attachment["name"]), small_style))
                        story.append(Spacer(1, 6))
                    except Exception:
                        story.append(
                            Paragraph(
                                f"Imagem nao renderizada: {pdf_text(attachment['name'])}",
                                small_style,
                            )
                        )

            story.append(Spacer(1, 8))

    story.append(PageBreak())
    story.append(Paragraph("Conclusao executiva", heading_style))
    story.append(
        Paragraph(
            f"Percentual de conformidade apurado: <b>{summary['rate']}%</b>. "
            f"{executive_reading(summary)}",
            body_style,
        )
    )
    doc.build(story, onFirstPage=draw_pdf_footer, onLaterPages=draw_pdf_footer)
    buffer.seek(0)
    return buffer.getvalue()


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(8.2), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(20, 54, 66)

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.62), Inches(1.1), Inches(8.1), Inches(0.4))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(13)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(77, 94, 102)


def add_metric_box(slide, left, top, title: str, value: str, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(1, left, top, Inches(2.0), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = value
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p2 = frame.add_paragraph()
    p2.text = title
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(10)
    p2.font.color.rgb = RGBColor(255, 255, 255)


def add_bullet(slide, text: str, top: float, bold_prefix: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.75), Inches(top), Inches(8.2), Inches(0.55))
    frame = box.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.font.size = Pt(15)
    paragraph.font.color.rgb = RGBColor(34, 42, 46)
    if bold_prefix and text.startswith(bold_prefix):
        run = paragraph.add_run()
        run.text = bold_prefix
        run.font.bold = True
        run.font.size = Pt(15)
        run.font.color.rgb = RGBColor(20, 54, 66)
        paragraph.add_run().text = text[len(bold_prefix):]
    else:
        paragraph.text = text


def make_pptx_report() -> bytes:
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5.625)
    blank_layout = presentation.slide_layouts[6]
    summary = calculate_summary()

    cover = presentation.slides.add_slide(blank_layout)
    add_title(cover, "Validacao de automacoes WhatsApp", "Relatorio executivo para diretoria")
    add_bullet(cover, f"Bateria: {st.session_state.audit_name or 'Nao informado'}", 2.0, "Bateria:")
    add_bullet(cover, f"Canal: {st.session_state.channel or 'Nao informado'}", 2.55, "Canal:")
    add_bullet(cover, f"Responsavel: {st.session_state.auditor or 'Nao informado'}", 3.1, "Responsavel:")
    add_bullet(cover, f"Data: {st.session_state.audit_date.strftime('%d/%m/%Y')}", 3.65, "Data:")

    metrics = presentation.slides.add_slide(blank_layout)
    add_title(metrics, "Resumo dos indicadores", executive_reading(summary))
    add_metric_box(metrics, Inches(0.7), Inches(2.0), "Total", str(summary["total"]), RGBColor(20, 54, 66))
    add_metric_box(metrics, Inches(2.95), Inches(2.0), "Conformes", str(summary["conform"]), RGBColor(27, 132, 87))
    add_metric_box(metrics, Inches(5.2), Inches(2.0), "Nao conformes", str(summary["non_conform"]), RGBColor(197, 76, 64))
    add_metric_box(metrics, Inches(7.45), Inches(2.0), "Conformidade", f'{summary["rate"]}%', RGBColor(38, 108, 164))

    detail = presentation.slides.add_slide(blank_layout)
    add_title(detail, "Principais resultados por teste")
    top = 1.45
    for index, test in enumerate(st.session_state.tests[:6], start=1):
        add_bullet(
            detail,
            f"{index}. {test['title']} - {status_label(test['status'])}",
            top,
        )
        top += 0.5
    if len(st.session_state.tests) > 6:
        add_bullet(detail, f"+ {len(st.session_state.tests) - 6} testes adicionais no relatorio PDF.", top)
    if not st.session_state.tests:
        add_bullet(detail, "Nenhum teste cadastrado ate o momento.", top)

    non_conform = [
        test for test in st.session_state.tests
        if test["status"] == "nao-conforme"
    ]
    action = presentation.slides.add_slide(blank_layout)
    add_title(action, "Pontos de atencao e plano de acao")
    if non_conform:
        top = 1.45
        for index, test in enumerate(non_conform[:5], start=1):
            note = test["notes"] or test["scenario"] or "Sem observacao registrada."
            add_bullet(action, f"{index}. {test['title']}: {note[:150]}", top)
            top += 0.68
    else:
        add_bullet(action, "Nao foram registrados testes nao conformes nesta bateria.", 1.6)
        add_bullet(action, "Recomendacao: manter monitoramento periodico e revalidar fluxos criticos.", 2.25)

    evidence_tests = [
        test for test in st.session_state.tests
        if test.get("attachments")
    ][:3]
    for test in evidence_tests:
        slide = presentation.slides.add_slide(blank_layout)
        add_title(slide, f"Evidencias - {test['title'][:45]}", status_label(test["status"]))
        left = 0.65
        top = 1.45
        for attachment in test.get("attachments", [])[:3]:
            try:
                slide.shapes.add_picture(
                    BytesIO(attachment["data"]),
                    Inches(left),
                    Inches(top),
                    width=Inches(2.7),
                )
                left += 3.05
            except Exception:
                add_bullet(slide, f"Imagem nao renderizada: {attachment['name']}", top)

    buffer = BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def reset_form() -> None:
    st.session_state.editing_id = None
    for key in ("title_input", "scenario_input", "expected_input", "notes_input"):
        st.session_state[key] = ""
    for index in range(st.session_state.get("sub_expected_count", 0)):
        st.session_state[f"sub_expected_input_{index}"] = ""
    st.session_state.sub_expected_count = 0
    st.session_state.status_input = "conforme"
    st.session_state.upload_version += 1


def add_sub_expected_field() -> None:
    st.session_state.sub_expected_count = st.session_state.get("sub_expected_count", 0) + 1


def current_sub_expected_values() -> list[str]:
    return [
        st.session_state.get(f"sub_expected_input_{index}", "").strip()
        for index in range(st.session_state.get("sub_expected_count", 0))
        if st.session_state.get(f"sub_expected_input_{index}", "").strip()
    ]


def uploaded_images() -> list[dict[str, object]]:
    key = f"attachments_input_{st.session_state.upload_version}"
    files = st.session_state.get(key) or []
    attachments = []

    for file in files:
        data = file.getvalue()
        size_mb = len(data) / (1024 * 1024)
        if size_mb > MAX_IMAGE_SIZE_MB:
            st.warning(
                f"{file.name} foi ignorado. Limite: {MAX_IMAGE_SIZE_MB} MB por imagem."
            )
            continue

        attachments.append(
            {
                "name": file.name,
                "type": file.type,
                "data": data,
            }
        )

    return attachments


def save_test() -> None:
    title = st.session_state.title_input.strip()
    if not title:
        st.warning("Informe o que será testado antes de salvar.")
        return

    payload = {
        "id": st.session_state.editing_id or str(uuid.uuid4()),
        "title": title,
        "scenario": st.session_state.scenario_input.strip(),
        "expected": st.session_state.expected_input.strip(),
        "sub_expected": current_sub_expected_values(),
        "notes": st.session_state.notes_input.strip(),
        "status": st.session_state.status_input,
        "attachments": uploaded_images(),
    }

    if st.session_state.editing_id:
        current = next(
            (item for item in st.session_state.tests if item["id"] == st.session_state.editing_id),
            None,
        )
        existing_attachments = current.get("attachments", []) if current else []
        payload["attachments"] = existing_attachments + payload["attachments"]
        st.session_state.tests = [
            payload if item["id"] == st.session_state.editing_id else item
            for item in st.session_state.tests
        ]
    else:
        st.session_state.tests.insert(0, payload)

    reset_form()
    st.success("Teste salvo.")


def edit_test(test_id: str) -> None:
    test = next((item for item in st.session_state.tests if item["id"] == test_id), None)
    if not test:
        return

    st.session_state.editing_id = test_id
    st.session_state.title_input = test["title"]
    st.session_state.scenario_input = test["scenario"]
    st.session_state.expected_input = test["expected"]
    sub_expected = test.get("sub_expected", [])
    st.session_state.sub_expected_count = len(sub_expected)
    for index, value in enumerate(sub_expected):
        st.session_state[f"sub_expected_input_{index}"] = value
    st.session_state.notes_input = test["notes"]
    st.session_state.status_input = test["status"]
    st.session_state.upload_version += 1


def delete_test(test_id: str) -> None:
    st.session_state.tests = [
        item for item in st.session_state.tests if item["id"] != test_id
    ]
    if st.session_state.editing_id == test_id:
        reset_form()


def render_header() -> None:
    left, right = st.columns([0.75, 0.25], vertical_alignment="center")
    with left:
        st.title("Validador de automações WhatsApp")
        st.caption("Registro de testes, conformidade e evidências do atendimento.")
    with right:
        if st.button("Sair", use_container_width=True):
            st.session_state.authenticated = False
            st.session_state.pending_password_change = False
            st.session_state.current_user = ""
            st.session_state.current_role = ""
            st.session_state.login_username = ""
            st.rerun()


def render_summary() -> None:
    summary = calculate_summary()
    cols = st.columns(4)
    cols[0].metric("Total de testes", summary["total"])
    cols[1].metric("Conformes", summary["conform"])
    cols[2].metric("Não conformes", summary["non_conform"])
    cols[3].metric("Conformidade", f'{summary["rate"]}%')


def render_audit_data() -> None:
    st.subheader("Dados da bateria")
    cols = st.columns(4)
    with cols[0]:
        st.text_input("Nome da bateria", key="audit_name")
    with cols[1]:
        st.text_input("Canal / número testado", key="channel")
    with cols[2]:
        st.text_input("Responsável", key="auditor")
    with cols[3]:
        st.date_input("Data", key="audit_date", format="DD/MM/YYYY")


def render_test_form() -> None:
    st.subheader("Novo teste")
    st.text_input(
        "O que vou testar",
        placeholder="Ex.: Segunda via de boleto",
        key="title_input",
    )

    cols = st.columns(2)
    with cols[0]:
        st.text_area(
            "Mensagem enviada / cenário",
            placeholder="Ex.: Cliente pergunta: Quero minha segunda via",
            key="scenario_input",
            height=120,
        )
        st.radio(
            "Resultado",
            options=["conforme", "nao-conforme", "pendente"],
            format_func=status_label,
            horizontal=True,
            key="status_input",
        )
    with cols[1]:
        st.text_area(
            "Resposta esperada da automação",
            placeholder="Ex.: Bot identifica o cliente e oferece link da segunda via",
            key="expected_input",
            height=120,
        )
        plus_cols = st.columns([0.86, 0.14])
        with plus_cols[1]:
            st.button(
                "+",
                key="add_sub_expected",
                help="Adicionar subcategoria ou sub resposta esperada",
                use_container_width=True,
                on_click=add_sub_expected_field,
            )
        for index in range(st.session_state.get("sub_expected_count", 0)):
            st.text_area(
                f"Subcategoria / sub resposta esperada {index + 1}",
                placeholder="Ex.: Se cliente não localizado, bot solicita CPF/CNPJ",
                key=f"sub_expected_input_{index}",
                height=90,
            )
        st.text_area(
            "Observações",
            placeholder="Evidências, horários, prints, falhas ou pontos de melhoria",
            key="notes_input",
            height=120,
        )

    upload_key = f"attachments_input_{st.session_state.upload_version}"
    files = st.file_uploader(
        "Evidências em imagem",
        type=["png", "jpg", "jpeg"],
        accept_multiple_files=True,
        key=upload_key,
        help="Aceita PNG, JPG e JPEG. Limite recomendado: 5 MB por imagem.",
    )

    if files:
        st.caption("Pré-visualização dos anexos")
        preview_cols = st.columns(min(3, len(files)))
        for index, file in enumerate(files):
            with preview_cols[index % len(preview_cols)]:
                st.image(file, caption=file.name, use_container_width=True)

    if st.session_state.editing_id:
        current = next(
            (item for item in st.session_state.tests if item["id"] == st.session_state.editing_id),
            None,
        )
        existing_attachments = current.get("attachments", []) if current else []
        if existing_attachments:
            st.caption("Anexos já salvos neste teste")
            preview_cols = st.columns(min(3, len(existing_attachments)))
            for index, attachment in enumerate(existing_attachments):
                with preview_cols[index % len(preview_cols)]:
                    st.image(
                        attachment["data"],
                        caption=attachment["name"],
                        use_container_width=True,
                    )

    label = "Salvar alteração" if st.session_state.editing_id else "Adicionar teste"
    cols = st.columns([0.25, 0.25, 0.5])
    with cols[0]:
        st.button(label, type="primary", use_container_width=True, on_click=save_test)
    with cols[1]:
        st.button("Cancelar edição", use_container_width=True, on_click=reset_form)


def render_report() -> None:
    st.subheader("Relatório final")
    if st.session_state.report_saved_message:
        st.success(st.session_state.report_saved_message)
        st.session_state.report_saved_message = ""
    if st.session_state.current_report_id:
        st.info("Você está editando um relatório salvo. Ao salvar, ele será atualizado sem duplicar.")

    summary = calculate_summary()
    st.write(
        f"**Bateria:** {st.session_state.audit_name or 'Não informado'}  \n"
        f"**Canal:** {st.session_state.channel or 'Não informado'}  \n"
        f"**Responsável:** {st.session_state.auditor or 'Não informado'}  \n"
        f"**Data:** {st.session_state.audit_date.strftime('%d/%m/%Y')}  \n"
        f"**Percentual de conformidade:** {summary['rate']}%"
    )

    csv_data = make_csv()
    st.download_button(
        "Baixar relatório CSV",
        data=csv_data.encode("utf-8-sig"),
        file_name=f"relatorio-testes-whatsapp-{st.session_state.audit_date}.csv",
        mime="text/csv",
        use_container_width=True,
    )
    col_pdf, col_ppt = st.columns(2)
    with col_pdf:
        st.download_button(
            "Baixar relatório executivo PDF",
            data=make_pdf_report(),
            file_name=f"{report_base_name()}-relatorio-executivo.pdf",
            mime="application/pdf",
            use_container_width=True,
            disabled=not st.session_state.tests,
        )
    with col_ppt:
        st.download_button(
            "Baixar modelo PowerPoint",
            data=make_pptx_report(),
            file_name=f"{report_base_name()}-apresentacao-diretoria.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
            disabled=not st.session_state.tests,
        )

    save_label = (
        "Atualizar relatório salvo"
        if st.session_state.current_report_id
        else "Salvar relatório no banco local"
    )
    st.button(
        save_label,
        type="primary",
        use_container_width=True,
        disabled=not st.session_state.tests,
        on_click=save_current_report_snapshot,
    )

    if not st.session_state.tests:
        st.info("Nenhum teste cadastrado ainda.")
        return

    for index, test in enumerate(st.session_state.tests, start=1):
        with st.container(border=True):
            cols = st.columns([0.65, 0.35], vertical_alignment="center")
            with cols[0]:
                st.markdown(f"**Teste {index}: {test['title']}**")
            with cols[1]:
                st.markdown(f"**Resultado:** {status_label(test['status'])}")

            st.write(f"**Cenário:** {test['scenario'] or 'Não informado.'}")
            st.write(f"**Esperado:** {test['expected'] or 'Não informado.'}")
            sub_expected = test.get("sub_expected", [])
            if sub_expected:
                st.write("**Sub respostas esperadas:**")
                for sub_index, sub_value in enumerate(sub_expected, start=1):
                    st.write(f"{sub_index}. {sub_value}")
            st.write(f"**Observações:** {test['notes'] or 'Não informado.'}")

            attachments = test.get("attachments", [])
            if attachments:
                st.write("**Evidências:**")
                image_cols = st.columns(min(3, len(attachments)))
                for image_index, attachment in enumerate(attachments):
                    with image_cols[image_index % len(image_cols)]:
                        st.image(
                            attachment["data"],
                            caption=attachment["name"],
                            use_container_width=True,
                        )

            action_cols = st.columns([0.18, 0.18, 0.64])
            with action_cols[0]:
                st.button(
                    "Editar",
                    key=f"edit_{test['id']}",
                    use_container_width=True,
                    on_click=edit_test,
                    args=(test["id"],),
                )
            with action_cols[1]:
                st.button(
                    "Excluir",
                    key=f"delete_{test['id']}",
                    use_container_width=True,
                    on_click=delete_test,
                    args=(test["id"],),
                )


def render_saved_reports() -> None:
    st.subheader("Banco local")
    status = database_status()
    status_cols = st.columns([0.2, 0.2, 0.2, 0.4])
    status_cols[0].metric("Status", "Instalado" if status["installed"] else "Pendente")
    status_cols[1].metric("Tipo", status["backend"])
    status_cols[2].metric("Relatórios salvos", status["reports"])
    status_cols[3].code(status["path"], language="text")

    if not status["persistent"]:
        st.warning(
            "Este app está usando SQLite local temporário. No Streamlit Cloud, usuários e relatórios "
            "podem ser perdidos quando o app reiniciar. Configure o secret DATABASE_URL com um "
            "PostgreSQL externo para deixar os dados permanentes."
        )

    if st.button("Instalar / verificar banco local", use_container_width=True):
        initialize_database()
        st.success("Banco local instalado e verificado com sucesso.")
        st.rerun()

    st.subheader("Relatórios salvos localmente")
    reports = list_saved_reports()

    if not reports:
        st.info("Nenhum relatório salvo no banco local.")
    else:
        st.caption(f"{len(reports)} relatório(s) salvo(s) para o usuário logado.")
        for report in reports:
            with st.container(border=True):
                created_at = datetime.fromisoformat(report["created_at"]).strftime("%d/%m/%Y %H:%M")
                st.markdown(
                    f"**{report['audit_name'] or 'Sem nome'}**  \n"
                    f"Canal: {report['channel'] or 'Não informado'}  \n"
                    f"Data do teste: {date.fromisoformat(report['audit_date']).strftime('%d/%m/%Y')}  \n"
                    f"Salvo em: {created_at}"
                )
                action_cols = st.columns(2)
                with action_cols[0]:
                    st.button(
                        "Carregar relatório",
                        key=f"load_saved_{report['id']}",
                        use_container_width=True,
                        on_click=load_saved_report,
                        args=(report["id"],),
                    )
                    st.download_button(
                        "Baixar PDF",
                        data=make_saved_report_pdf(report),
                        file_name=(
                            f"{safe_report_base_name(report['audit_name'])}-"
                            f"{report['audit_date']}-relatorio-executivo.pdf"
                        ),
                        mime="application/pdf",
                        key=f"download_saved_pdf_{report['id']}",
                        use_container_width=True,
                    )
                with action_cols[1]:
                    with st.form(f"delete_saved_form_{report['id']}"):
                        delete_password = st.text_input(
                            "Senha para excluir este relatório",
                            type="password",
                            key=f"delete_saved_password_{report['id']}",
                        )
                        delete_submitted = st.form_submit_button(
                            "Excluir relatório",
                            use_container_width=True,
                        )

                    if delete_submitted:
                        if verify_current_password(delete_password):
                            delete_saved_report(report["id"])
                            st.rerun()
                        else:
                            st.error("Senha inválida. Relatório não foi apagado.")

    with st.form("clear_saved_reports_form"):
        st.warning("Para limpar os relatórios salvos, confirme sua senha atual.")
        password = st.text_input("Senha para limpar relatórios salvos", type="password")
        submitted = st.form_submit_button("Limpar relatórios salvos", use_container_width=True)

    if submitted:
        if verify_current_password(password):
            clear_saved_reports()
            st.rerun()
        else:
            st.error("Senha inválida. Nada foi apagado.")


def render_admin_panel() -> None:
    if st.session_state.get("current_role") != "admin":
        return

    st.divider()
    st.subheader("Painel administrador")
    users = list_all_users()
    reports = list_all_saved_reports()

    cols = st.columns(3)
    cols[0].metric("Usuários", len(users))
    cols[1].metric("Relatórios no banco", len(reports))
    cols[2].metric("Administradores", sum(1 for user in users if user["role"] == "admin"))

    st.markdown("**Usuários cadastrados**")
    if users:
        user_rows = [
            {
                "Usuário": user["username"],
                "Perfil": user["role"],
                "Troca pendente": "Sim" if user["must_change_password"] else "Não",
                "Relatórios": user["saved_reports"],
                "Último login": user["last_login_at"] or "Nunca",
            }
            for user in users
        ]
        st.dataframe(user_rows, use_container_width=True, hide_index=True)
    else:
        st.info("Nenhum usuário cadastrado.")

    st.markdown("**Gerenciar acessos**")
    manage_cols = st.columns(3)
    with manage_cols[0]:
        with st.form("admin_create_user_form"):
            st.markdown("Novo usuário provisório")
            new_user = st.text_input("Usuário", key="admin_create_username")
            new_user_password = st.text_input(
                "Senha provisória",
                type="password",
                key="admin_create_password",
            )
            new_user_role = st.selectbox(
                "Perfil",
                options=["user", "admin"],
                format_func=lambda value: "Administrador" if value == "admin" else "Usuário",
                key="admin_create_role",
            )
            admin_password = st.text_input(
                "Senha do administrador",
                type="password",
                key="admin_create_admin_password",
            )
            create_submitted = st.form_submit_button("Criar usuário", use_container_width=True)

        if create_submitted:
            normalized_new_user = normalize_username(new_user)
            if not verify_current_password(admin_password):
                st.error("Senha admin inválida. Usuário não foi criado.")
            elif len(normalized_new_user) < 3:
                st.error("O usuário precisa ter pelo menos 3 caracteres.")
            elif len(new_user_password) < 8:
                st.error("A senha provisória precisa ter pelo menos 8 caracteres.")
            elif get_user(normalized_new_user):
                st.error("Este usuário já existe.")
            else:
                create_user(normalized_new_user, new_user_password, new_user_role, True)
                st.success("Usuário criado. No primeiro login ele deverá escolher novo usuário e senha.")
                st.rerun()

    with manage_cols[1]:
        user_options = [user["username"] for user in users]
        with st.form("admin_reset_user_form"):
            st.markdown("Resetar senha de usuário")
            selected_user = st.selectbox(
                "Usuário",
                options=user_options,
                key="admin_reset_username",
            )
            temporary_password = st.text_input(
                "Nova senha provisória",
                type="password",
                key="admin_reset_password",
            )
            reset_admin_password = st.text_input(
                "Senha do administrador",
                type="password",
                key="admin_reset_admin_password",
            )
            reset_submitted = st.form_submit_button("Resetar senha", use_container_width=True)

        if reset_submitted:
            if not verify_current_password(reset_admin_password):
                st.error("Senha admin inválida. Senha não foi resetada.")
            elif not selected_user:
                st.error("Selecione um usuário.")
            elif len(temporary_password) < 8:
                st.error("A senha provisória precisa ter pelo menos 8 caracteres.")
            else:
                reset_user_password(selected_user, temporary_password)
                st.success("Senha resetada. No próximo login o usuário deverá escolher novo usuário e senha.")
                st.rerun()

    with manage_cols[2]:
        user_options = [user["username"] for user in users]
        with st.form("admin_delete_user_form"):
            st.markdown("Excluir usuário")
            delete_username = st.selectbox(
                "Usuário",
                options=user_options,
                key="admin_delete_username",
            )
            delete_user_reports = next(
                (user["saved_reports"] for user in users if user["username"] == delete_username),
                0,
            )
            st.caption(f"Também serão excluídos {delete_user_reports} relatório(s) deste usuário.")
            delete_admin_password = st.text_input(
                "Senha do administrador",
                type="password",
                key="admin_delete_user_admin_password",
            )
            delete_submitted = st.form_submit_button("Excluir usuário", use_container_width=True)

        if delete_submitted:
            selected_user_data = next(
                (user for user in users if user["username"] == delete_username),
                None,
            )
            admin_count = sum(1 for user in users if user["role"] == "admin")
            if not verify_current_password(delete_admin_password):
                st.error("Senha admin inválida. Usuário não foi excluído.")
            elif not selected_user_data:
                st.error("Selecione um usuário válido.")
            elif delete_username == st.session_state.get("current_user"):
                st.error("Você não pode excluir o usuário que está logado agora.")
            elif selected_user_data["role"] == "admin" and admin_count <= 1:
                st.error("Não é possível excluir o último administrador.")
            else:
                delete_user(delete_username)
                st.success("Usuário e relatórios vinculados foram excluídos.")
                st.rerun()

    st.markdown("**Relatórios de todos os usuários**")
    if not reports:
        st.info("Nenhum relatório salvo por usuários.")
        return

    for report in reports:
        with st.container(border=True):
            created_at = datetime.fromisoformat(report["created_at"]).strftime("%d/%m/%Y %H:%M")
            st.markdown(
                f"**{report['audit_name'] or 'Sem nome'}**  \n"
                f"Usuário: {report['username']}  \n"
                f"Canal: {report['channel'] or 'Não informado'}  \n"
                f"Data do teste: {date.fromisoformat(report['audit_date']).strftime('%d/%m/%Y')}  \n"
                f"Salvo em: {created_at}"
            )
            action_cols = st.columns(2)
            with action_cols[0]:
                st.button(
                    "Carregar relatório",
                    key=f"admin_load_saved_{report['id']}",
                    use_container_width=True,
                    on_click=load_any_saved_report,
                    args=(report["id"],),
                )
                st.download_button(
                    "Baixar PDF",
                    data=make_saved_report_pdf(report),
                    file_name=(
                        f"{safe_report_base_name(report['audit_name'])}-"
                        f"{report['audit_date']}-relatorio-executivo.pdf"
                    ),
                    mime="application/pdf",
                    key=f"admin_download_saved_pdf_{report['id']}",
                    use_container_width=True,
                )
            with action_cols[1]:
                with st.form(f"admin_delete_saved_form_{report['id']}"):
                    admin_password = st.text_input(
                        "Senha admin para excluir",
                        type="password",
                        key=f"admin_delete_saved_password_{report['id']}",
                    )
                    submitted = st.form_submit_button(
                        "Excluir relatório",
                        use_container_width=True,
                    )

                if submitted:
                    if verify_current_password(admin_password):
                        delete_any_saved_report(report["id"])
                        st.rerun()
                    else:
                        st.error("Senha admin inválida. Relatório não foi apagado.")


def main() -> None:
    init_state()

    if st.session_state.pending_password_change:
        password_change_screen()
        return

    if not st.session_state.authenticated:
        login_screen()
        return

    render_header()
    render_summary()
    st.divider()

    render_audit_data()
    st.divider()

    left, right = st.columns([0.52, 0.48], gap="large")
    with left:
        render_test_form()
    with right:
        render_report()

    st.divider()
    render_saved_reports()
    render_admin_panel()


if __name__ == "__main__":
    main()
